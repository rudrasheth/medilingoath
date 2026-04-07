from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # --- Helper Function to Add Title & Content Slides ---
    def add_slide(prs, title_text, content_text_list, image_placeholder_text=None):
        slide_layout = prs.slide_layouts[1] # Using 'Title and Content' layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set Title
        title = slide.shapes.title
        title.text = title_text
        
        # Format Title (Optional: Make it bold/blue to match tech theme)
        for paragraph in title.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(36)
            paragraph.font.color.rgb = RGBColor(0, 51, 102) # Dark Blue

        # Set Content (Bullet Points)
        shapes = slide.shapes
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.text = content_text_list[0] # First bullet
        
        for item in content_text_list[1:]:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(20)

        # Add Image Placeholder (Right side)
        if image_placeholder_text:
            left = Inches(5.5)
            top = Inches(2.0)
            width = Inches(4.0)
            height = Inches(4.0)
            shape = slide.shapes.add_shape(
                1, left, top, width, height # 1 is the code for a Rectangle
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(220, 220, 220) # Grey placeholder
            shape.text = f"[PLACEHOLDER]\n\n{image_placeholder_text}\n(Paste Image Here)"
            
            # Adjust text box on the left to make room for image
            body_shape.width = Inches(4.5)

    # =======================================================
    # SLIDE 1: TITLE SLIDE
    # =======================================================
    slide_layout = prs.slide_layouts[0] # Title Slide Layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Integrated Community Resource & Crisis Response Platform"
    subtitle.text = "Team Name: LLMao\nTeam Leader: Yati Rathod\nCollege: Dwarkadas J. Sanghvi College of Engineering\nTheme: Space / Crisis Response"

    # =======================================================
    # SLIDE 2: PROBLEM STATEMENT
    # =======================================================
    # Source [16-23], [59-60]
    problem_points = [
        "Fragmented Coordination: Disconnect between on-ground info, resources, and agencies.",
        "Communication Blackouts: Standard apps fail when cellular towers go down during disasters.",
        "The Visibility Gap: Victims cannot report location; Agencies cannot track real-time inventory.",
        "Unverified Noise: Panic leads to fake reports, overwhelming 100/108 lines."
    ]
    add_slide(prs, "PROBLEM STATEMENT", problem_points, "Insert Image: Collage of 'No Network' signal + Chaos/Traffic Map")

    # =======================================================
    # SLIDE 3: SOLUTION (USP)
    # =======================================================
    # Source [28-34], [65-66]
    solution_points = [
        "Offline-First Mesh (USP): Bluetooth 'hopping' allows reporting without Internet.",
        "Agentic AI Support (USP): Voice bot talks to users via phone call, extracting location & severity automatically.",
        "Unified Command Dashboard: Real-time heatmaps for agencies to visualize clusters.",
        "Smart Allocation: Algorithms match 'Demand' (Victims) with nearest 'Supply' (Volunteers)."
    ]
    add_slide(prs, "PROPOSED SOLUTION & USP", solution_points, "Insert Image: System Architecture Diagram (App <-> Mesh <-> Cloud)")

    # =======================================================
    # SLIDE 4: TECH STACK
    # =======================================================
    # Source [71-72]
    tech_points = [
        "Frontend: React Native (Mobile PWA), React.js (Admin Dashboard)",
        "Connectivity: Bluetooth LE (Mesh), Twilio SDK (SMS/Voice)",
        "Backend: Node.js (Real-time events), Python (AI Processing)",
        "Data & ML: PostgreSQL + PostGIS (Location), OpenAI/Llama (Voice Intelligence)"
    ]
    add_slide(prs, "TECH STACK", tech_points, "Insert Image: Tech Logos (React, Twilio, Node, Postgres)")

    # =======================================================
    # SLIDE 5: IMPACT AND BENEFITS
    # =======================================================
    # Source [78-80], [44-46]
    impact_points = [
        "100% Inclusivity: Works via Smartphone, SMS, or Voice Call.",
        "Zero-Network Resilience: Communication continues even during total blackouts via Mesh.",
        "Faster Response: Automated triage reduces decision time by filtering noise.",
        "Responder Safety: 'Safe Routing' guides volunteers away from active hazards."
    ]
    add_slide(prs, "IMPACT AND BENEFITS", impact_points, "Insert Image: Before vs After Scenario (Chaos vs Organized Pins)")

    # =======================================================
    # SLIDE 6: FREE PAGE (Business & Future)
    # =======================================================
    # Source [86-88]
    future_points = [
        "Business Model: B2G Subscription for Disaster Mgmt Authorities.",
        "Drone Integration: Autonomous delivery of medical kits to 'Red Zones'.",
        "Predictive Analytics: AI predicts resource demand based on historical weather patterns.",
        "Global Ledger: Blockchain transparency for relief fund tracking."
    ]
    add_slide(prs, "FUTURE SCOPE & BUSINESS MODEL", future_points, "Insert Image: Roadmap Timeline or Drone Concept Art")

    # --- Save Presentation ---
    prs.save('Team_LLMao_Submission.pptx')
    print("Presentation 'Team_LLMao_Submission.pptx' created successfully!")

if __name__ == "__main__":
    create_presentation()